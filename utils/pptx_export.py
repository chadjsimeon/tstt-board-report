"""
TSTT Board Report — PowerPoint export module.

Entry point:  build_pptx(data, ar_data=None) -> io.BytesIO

Dependencies: python-pptx>=1.0.2, kaleido==0.2.1
"""

from __future__ import annotations

import io
import math
from typing import Optional

import pandas as pd
import plotly.graph_objects as go
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

import utils.slide_charts as sc
from utils.data_loader import get_month_order

# ── Geometry ──────────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

MARGIN_L    = Inches(0.30)
MARGIN_R    = Inches(0.30)
HEADER_TOP  = Inches(0.10)
HEADER_H    = Inches(0.60)
CONTENT_TOP = Inches(0.80)
CONTENT_H   = Inches(6.45)
FOOTER_TOP  = Inches(7.28)
FOOTER_H    = Inches(0.18)
USABLE_W    = SLIDE_W - MARGIN_L - MARGIN_R   # 12.73 in
GAP         = Inches(0.18)

KPI_H       = Inches(0.95)
CHART_TOP_1 = CONTENT_TOP + KPI_H + GAP          # below a KPI row
CHART_H_1   = CONTENT_H - KPI_H - GAP            # height when KPI row above
CHART_H_FULL = CONTENT_H                          # full content height, no KPI row

# Pre-calculated widths
HALF_W  = (USABLE_W - GAP)      / 2
THIRD_W = (USABLE_W - 2 * GAP)  / 3
W4      = (USABLE_W - 3 * GAP)  / 4
W3      = (USABLE_W - 2 * GAP)  / 3

# ── Colour palette ────────────────────────────────────────────────────────────

BG_COLOR   = RGBColor(0x00, 0x00, 0x00)
HDR_COLOR  = RGBColor(0x0f, 0x1e, 0x3c)
CARD_COLOR = RGBColor(0x1a, 0x1a, 0x2e)
BORDER_CLR = RGBColor(0x2a, 0x2a, 0x4a)
GREEN_RGB  = RGBColor(0x00, 0xff, 0x88)
BLUE_RGB   = RGBColor(0x44, 0x88, 0xff)
PURPLE_RGB = RGBColor(0xaa, 0x44, 0xff)
YELLOW_RGB = RGBColor(0xff, 0xd7, 0x00)
ORANGE_RGB = RGBColor(0xff, 0x88, 0x44)
WHITE_RGB  = RGBColor(0xff, 0xff, 0xff)
MUTED_RGB  = RGBColor(0x77, 0x88, 0xaa)
RED_RGB    = RGBColor(0xef, 0x44, 0x44)

# Section accent colours
SECTION_ACCENTS = {
    "Financial": BLUE_RGB,
    "Customer":  GREEN_RGB,
    "Network":   PURPLE_RGB,
    "People":    YELLOW_RGB,
}

# ── Kaleido render sizes (width × height in px, rendered at scale=2) ──────────
PX_FULL_W, PX_FULL_H   = 2400, 800
PX_HALF_W, PX_HALF_H   = 1180, 800
PX_TALL_W, PX_TALL_H   = 2400, 960
PX_THIRD_W, PX_THIRD_H = 780,  800


# ── Low-level helpers ─────────────────────────────────────────────────────────

def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fill_bg(slide, color: RGBColor = BG_COLOR) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, left, top, width, height,
          fill: RGBColor = CARD_COLOR, line: Optional[RGBColor] = BORDER_CLR,
          line_w: float = 0.75):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def _txb(slide, text: str, left, top, width, height,
         size: float = 12, bold: bool = False, color: RGBColor = WHITE_RGB,
         align: PP_ALIGN = PP_ALIGN.LEFT, wrap: bool = True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def _fig_to_png(fig: go.Figure, w_px: int, h_px: int) -> bytes:
    fig.update_layout(paper_bgcolor="rgba(0,0,0,0)",
                      plot_bgcolor="rgba(0,0,0,0)")
    try:
        return fig.to_image(format="png", width=w_px, height=h_px, scale=2)
    except Exception as exc:
        raise RuntimeError(
            f"Kaleido render failed. Ensure kaleido>=1.0.0 is installed. Error: {exc}"
        ) from exc


def _add_chart(slide, fig: go.Figure, left, top, width, height,
               px_w: int = PX_FULL_W, px_h: int = PX_FULL_H):
    if fig is None or (hasattr(fig, "data") and len(fig.data) == 0):
        return
    png = _fig_to_png(fig, px_w, px_h)
    slide.shapes.add_picture(io.BytesIO(png), left, top, width, height)


def _add_footer(slide, text: str = "CONFIDENTIAL — Board Use Only"):
    _txb(slide, text,
         left=MARGIN_L, top=FOOTER_TOP,
         width=USABLE_W, height=FOOTER_H,
         size=7, color=MUTED_RGB, align=PP_ALIGN.CENTER)


# ── Header ────────────────────────────────────────────────────────────────────

def _add_header(slide, title: str, subtitle: str = "", period: str = "") -> None:
    # Dark navy strip
    _rect(slide, left=0, top=HEADER_TOP, width=SLIDE_W, height=HEADER_H,
          fill=HDR_COLOR, line=None)

    # "TSTT" logo mark — green box
    logo_w, logo_h = Inches(0.68), Inches(0.38)
    logo_l = MARGIN_L
    logo_t = HEADER_TOP + (HEADER_H - logo_h) / 2
    logo_box = _rect(slide, logo_l, logo_t, logo_w, logo_h,
                     fill=RGBColor(0x00, 0x11, 0x22), line=GREEN_RGB, line_w=1.5)
    tf = logo_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "TSTT"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = GREEN_RGB

    # Title text
    title_l = logo_l + logo_w + Inches(0.16)
    title_w = USABLE_W - logo_w - Inches(0.16) - (Inches(2.2) if period else 0)
    _txb(slide, title, left=title_l, top=HEADER_TOP + Inches(0.10),
         width=title_w, height=Inches(0.28),
         size=14, bold=True, color=WHITE_RGB)
    if subtitle:
        _txb(slide, subtitle,
             left=title_l, top=HEADER_TOP + Inches(0.38),
             width=title_w, height=Inches(0.18),
             size=9, color=MUTED_RGB)

    # Period badge (right side)
    if period:
        badge_w = Inches(2.0)
        badge_l = SLIDE_W - MARGIN_R - badge_w
        badge_t = HEADER_TOP + (HEADER_H - Inches(0.32)) / 2
        badge_box = _rect(slide, badge_l, badge_t, badge_w, Inches(0.32),
                          fill=RGBColor(0x00, 0x22, 0x11), line=GREEN_RGB, line_w=1)
        tf = badge_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"YTD {period}"
        run.font.size = Pt(9)
        run.font.color.rgb = GREEN_RGB


# ── KPI card ──────────────────────────────────────────────────────────────────

def _add_kpi_box(slide, label: str, value: str, sub: str,
                 left, top, width, height=KPI_H,
                 accent: RGBColor = BLUE_RGB) -> None:
    # Card body
    _rect(slide, left, top, width, height, fill=CARD_COLOR, line=BORDER_CLR)

    # Accent top bar (3px tall rectangle at very top of card)
    accent_bar_h = Pt(3)
    _rect(slide, left, top, width, accent_bar_h, fill=accent, line=None)

    inner_l = left + Inches(0.15)
    inner_w = width - Inches(0.30)

    # Label
    _txb(slide, label.upper(),
         left=inner_l, top=top + Inches(0.12),
         width=inner_w, height=Inches(0.20),
         size=7.5, color=MUTED_RGB)

    # Value
    _txb(slide, value,
         left=inner_l, top=top + Inches(0.31),
         width=inner_w, height=Inches(0.38),
         size=18, bold=True, color=WHITE_RGB)

    # Sub-text
    if sub:
        _txb(slide, sub,
             left=inner_l, top=top + Inches(0.68),
             width=inner_w, height=Inches(0.20),
             size=8, color=accent)


def _kpi_row_4(slide, kpis: list[tuple], top=CONTENT_TOP) -> None:
    """Draw 4 KPI boxes across the full usable width."""
    for i, (label, value, sub, accent) in enumerate(kpis):
        left = MARGIN_L + i * (W4 + GAP)
        _add_kpi_box(slide, label, value, sub, left, top, W4, KPI_H, accent)


def _kpi_row_3(slide, kpis: list[tuple], top=CONTENT_TOP) -> None:
    for i, (label, value, sub, accent) in enumerate(kpis):
        left = MARGIN_L + i * (W3 + GAP)
        _add_kpi_box(slide, label, value, sub, left, top, W3, KPI_H, accent)


# ── Utility: get period string ────────────────────────────────────────────────

def _period(data: dict) -> str:
    try:
        fin = data.get("Financial_Monthly", pd.DataFrame())
        if not fin.empty:
            return pd.to_datetime(fin["Month"].iloc[-1],
                                  format="%b-%y").strftime("%B %Y")
    except Exception:
        pass
    try:
        cc = data.get("Cash_CAPEX", pd.DataFrame())
        if not cc.empty:
            return pd.to_datetime(cc["Month"].iloc[-1],
                                  format="%b-%y").strftime("%B %Y")
    except Exception:
        pass
    return ""


def _fmt(val, prefix="", suffix="M", decimals=0) -> str:
    try:
        return f"{prefix}{abs(float(val)):,.{decimals}f}{suffix}"
    except (TypeError, ValueError):
        return "—"


def _safe_sum(df, col):
    if col in df.columns:
        return float(df[col].sum())
    return 0.0


def _safe_last(df, col):
    if col in df.columns:
        v = df[col].dropna()
        return float(v.iloc[-1]) if len(v) else None
    return None


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def _slide_title(prs: Presentation, period: str) -> None:
    sl = _blank_slide(prs)
    _fill_bg(sl, BG_COLOR)

    # Vertical centred block
    block_top = Inches(2.2)

    # Logo mark — large, centred
    logo_w, logo_h = Inches(1.6), Inches(0.75)
    logo_l = (SLIDE_W - logo_w) / 2
    logo_box = _rect(sl, logo_l, block_top, logo_w, logo_h,
                     fill=RGBColor(0x00, 0x11, 0x22), line=GREEN_RGB, line_w=2)
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "TSTT"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = GREEN_RGB

    # Report title
    _txb(sl, "Board Report",
         left=0, top=block_top + Inches(0.95), width=SLIDE_W, height=Inches(0.85),
         size=40, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER)

    # Period
    if period:
        _txb(sl, period,
             left=0, top=block_top + Inches(1.90), width=SLIDE_W, height=Inches(0.45),
             size=18, color=GREEN_RGB, align=PP_ALIGN.CENTER)

    # Divider line
    div = sl.shapes.add_shape(1,
        left=Inches(4.5), top=block_top + Inches(2.48),
        width=Inches(4.33), height=Pt(1))
    div.fill.solid()
    div.fill.fore_color.rgb = RGBColor(0x2a, 0x3a, 0x5a)
    div.line.fill.background()

    # Confidential
    _txb(sl, "CONFIDENTIAL — Board Members Only",
         left=0, top=block_top + Inches(2.65), width=SLIDE_W, height=Inches(0.35),
         size=10, color=MUTED_RGB, align=PP_ALIGN.CENTER)


def _slide_scorecard(prs: Presentation, data: dict, period: str) -> None:
    kpi = data.get("KPI_Summary", pd.DataFrame())
    sl = _blank_slide(prs)
    _fill_bg(sl)
    _add_header(sl, "Enterprise KPI Scorecard", "", period)
    _add_footer(sl)

    if kpi.empty:
        _txb(sl, "KPI data not available",
             MARGIN_L, CONTENT_TOP, USABLE_W, CONTENT_H,
             size=14, color=MUTED_RGB, align=PP_ALIGN.CENTER)
        return

    latest_month = kpi["Month"].iloc[-1]
    latest_kpi   = kpi[kpi["Month"] == latest_month].copy()
    sections     = latest_kpi["Section"].unique().tolist()

    # Build rows: group by section
    all_rows = []
    for sec in sections:
        rows = latest_kpi[latest_kpi["Section"] == sec]
        all_rows.append(("SECTION", sec, "", "", ""))
        for _, r in rows.iterrows():
            unit   = str(r.get("Unit", ""))
            actual = r.get("Actual")
            aop    = r.get("AOP")
            status = str(r.get("Status", ""))
            try:
                av = float(actual)
                val_str = (f"{av:,.1f}" if unit == ""
                           else f"{av:.1f}%" if unit == "%" else f"{av:,.0f}")
            except (TypeError, ValueError):
                val_str = "—"
            try:
                aop_pct = (float(actual) - float(aop)) / float(aop) * 100
                vs_str  = f"{aop_pct:+.1f}%"
            except (TypeError, ValueError, ZeroDivisionError):
                vs_str = "—"
            all_rows.append(("KPI", str(r["KPI_Name"]), val_str, vs_str, status))

    n_rows = len(all_rows)
    table_h = CONTENT_H - Inches(0.05)
    row_h = table_h / max(n_rows, 1)

    # Draw rows manually as shapes + text (python-pptx tables are complex to style dark)
    col_xs = [MARGIN_L,
               MARGIN_L + Inches(5.5),
               MARGIN_L + Inches(8.0),
               MARGIN_L + Inches(9.8)]
    col_ws = [Inches(5.2), Inches(2.3), Inches(1.7), Inches(3.0)]

    for ri, (rtype, c0, c1, c2, c3) in enumerate(all_rows):
        row_top = CONTENT_TOP + ri * row_h

        if rtype == "SECTION":
            # Section header bar
            _rect(sl, MARGIN_L, row_top, USABLE_W, row_h,
                  fill=HDR_COLOR, line=None)
            accent = SECTION_ACCENTS.get(c0, BLUE_RGB)
            _rect(sl, MARGIN_L, row_top, Pt(3), row_h, fill=accent, line=None)
            _txb(sl, c0.upper(),
                 left=MARGIN_L + Inches(0.10), top=row_top,
                 width=USABLE_W, height=row_h,
                 size=9, bold=True, color=accent)
        else:
            # Alternating row bg
            bg = RGBColor(0x12, 0x12, 0x22) if ri % 2 == 0 else BG_COLOR
            _rect(sl, MARGIN_L, row_top, USABLE_W, row_h, fill=bg, line=None)

            # KPI Name
            _txb(sl, c0,
                 left=col_xs[0] + Inches(0.12), top=row_top,
                 width=col_ws[0], height=row_h,
                 size=9, color=WHITE_RGB)

            # Actual
            _txb(sl, c1,
                 left=col_xs[1], top=row_top,
                 width=col_ws[1], height=row_h,
                 size=9, bold=True, color=WHITE_RGB)

            # vs AOP
            if c2 and c2 != "—":
                try:
                    pct = float(c2.replace("%", "").replace("+", ""))
                    aop_col = GREEN_RGB if pct >= 0 else RED_RGB
                except ValueError:
                    aop_col = MUTED_RGB
            else:
                aop_col = MUTED_RGB
            _txb(sl, c2,
                 left=col_xs[2], top=row_top,
                 width=col_ws[2], height=row_h,
                 size=9, color=aop_col)

            # Status
            status_map = {"🟢": "●", "🔴": "●", "🟡": "●"}
            color_map  = {"🟢": GREEN_RGB, "🔴": RED_RGB, "🟡": YELLOW_RGB}
            dot = status_map.get(c3, c3)
            dot_col = color_map.get(c3, MUTED_RGB)
            _txb(sl, dot,
                 left=col_xs[3], top=row_top,
                 width=col_ws[3], height=row_h,
                 size=11, color=dot_col)


def _slide_financial(prs: Presentation, data: dict, period: str) -> None:
    fin = data.get("Financial_Monthly", pd.DataFrame())
    sl = _blank_slide(prs)
    _fill_bg(sl)
    _add_header(sl, "Financial Performance",
                "Revenue · EBITDA · PAT · Margin", period)
    _add_footer(sl)

    rev_ytd    = _safe_sum(fin, "Revenue")
    ebitda_ytd = _safe_sum(fin, "EBITDA")
    pat_ytd    = _safe_sum(fin, "PAT")
    margin_val = _safe_last(fin, "EBITDA_Margin")
    rev_aop    = _safe_sum(fin, "Revenue_AOP")
    ebitda_aop = _safe_sum(fin, "EBITDA_AOP")

    def _vs(act, plan):
        try:
            return f"{(act - plan) / plan * 100:+.1f}% vs AOP"
        except (ZeroDivisionError, TypeError):
            return "—"

    kpis = [
        ("Revenue YTD",    _fmt(rev_ytd),    _vs(rev_ytd, rev_aop),    BLUE_RGB),
        ("EBITDA YTD",     _fmt(ebitda_ytd), _vs(ebitda_ytd, ebitda_aop), GREEN_RGB),
        ("PAT YTD",        _fmt(pat_ytd),    "",                        PURPLE_RGB),
        ("EBITDA Margin",  f"{margin_val:.1f}%" if margin_val else "—", "Latest month", YELLOW_RGB),
    ]
    _kpi_row_4(sl, kpis, top=CONTENT_TOP)

    fig_rev = sc.fig_revenue_trend(data, height=440)
    fig_eb  = sc.fig_ebitda_trend(data, height=440)
    chart_top = CHART_TOP_1
    chart_h   = CHART_H_1

    _add_chart(sl, fig_rev,
               left=MARGIN_L, top=chart_top,
               width=HALF_W, height=chart_h,
               px_w=PX_HALF_W, px_h=PX_HALF_H)
    _add_chart(sl, fig_eb,
               left=MARGIN_L + HALF_W + GAP, top=chart_top,
               width=HALF_W, height=chart_h,
               px_w=PX_HALF_W, px_h=PX_HALF_H)


def _slide_ebitda_bridge(prs: Presentation, data: dict, period: str) -> None:
    bridge = data.get("EBITDA_Bridge", pd.DataFrame())
    sl = _blank_slide(prs)
    _fill_bg(sl)
    _add_header(sl, "EBITDA Bridge", "AOP vs Actual — variance analysis", period)
    _add_footer(sl)

    aop_val    = 0.0
    actual_val = 0.0
    if not bridge.empty:
        aop_row = bridge[bridge["Type"] == "Absolute"]
        tot_row = bridge[bridge["Type"] == "Total"]
        if not aop_row.empty:
            aop_val = float(aop_row["Value"].iloc[0])
        if not tot_row.empty:
            actual_val = float(tot_row["Value"].iloc[-1])

    gap_val = actual_val - aop_val

    kpis = [
        ("AOP EBITDA",    _fmt(aop_val),    "",        BLUE_RGB),
        ("Actual EBITDA", _fmt(actual_val), "",        GREEN_RGB if actual_val >= aop_val else RED_RGB),
        ("Gap to AOP",    _fmt(gap_val, prefix="" + ("+" if gap_val >= 0 else "")),
         "positive = ahead", GREEN_RGB if gap_val >= 0 else RED_RGB),
    ]
    _kpi_row_3(sl, kpis, top=CONTENT_TOP)

    fig = sc.fig_ebitda_bridge(data, height=500)
    _add_chart(sl, fig,
               left=MARGIN_L, top=CHART_TOP_1,
               width=USABLE_W, height=CHART_H_1,
               px_w=PX_TALL_W, px_h=PX_TALL_H)


def _slide_opex(prs: Presentation, data: dict, period: str) -> None:
    opex = data.get("OPEX", pd.DataFrame())
    sl = _blank_slide(prs)
    _fill_bg(sl)
    _add_header(sl, "OPEX & Cost Management",
                "Actual vs Plan by category — latest month", period)
    _add_footer(sl)

    if not opex.empty and "Month" in opex.columns:
        latest_month = opex["Month"].iloc[-1]
        m_data = opex[opex["Month"] == latest_month]
        total_act  = _safe_sum(m_data, "Actual")
        total_plan = _safe_sum(m_data, "Plan")
        variance   = total_act - total_plan
    else:
        total_act = total_plan = variance = 0.0

    kpis = [
        ("Total OPEX Actual",  _fmt(total_act),  "",    BLUE_RGB),
        ("Total OPEX Plan",    _fmt(total_plan),  "",   MUTED_RGB),
        ("Variance",           _fmt(variance, prefix="" + ("+" if variance >= 0 else "")),
         "over budget" if variance > 0 else "under budget",
         RED_RGB if variance > 0 else GREEN_RGB),
    ]
    _kpi_row_3(sl, kpis, top=CONTENT_TOP)

    fig = sc.fig_opex_bullet(data, height=500)
    _add_chart(sl, fig,
               left=MARGIN_L, top=CHART_TOP_1,
               width=USABLE_W, height=CHART_H_1,
               px_w=PX_TALL_W, px_h=PX_TALL_H)


def _slide_cash_capex(prs: Presentation, data: dict,
                      ar_data: Optional[dict], period: str) -> None:
    cc = data.get("Cash_CAPEX", pd.DataFrame())
    sl = _blank_slide(prs)
    _fill_bg(sl)
    _add_header(sl, "Cash, Working Capital & CAPEX",
                "Cash balance · AR Aging · CAPEX spend", period)
    _add_footer(sl)

    latest = cc.iloc[-1] if not cc.empty else pd.Series(dtype=float)
    first  = cc.iloc[0]  if not cc.empty else pd.Series(dtype=float)

    def _safe(v, d=0.0):
        try: return float(v) if pd.notna(v) else d
        except (TypeError, ValueError): return d

    cash_val   = _safe(latest.get("Cash_Balance", 0))
    debt_val   = _safe(latest.get("Net_Debt", 0))
    fcf_ytd    = float(cc["FCF"].sum()) if "FCF" in cc.columns else 0.0
    capex_act  = float(cc["CAPEX_Actual"].sum()) if "CAPEX_Actual" in cc.columns else 0.0

    fcf_col = GREEN_RGB if fcf_ytd >= 0 else RED_RGB
    kpis = [
        ("Cash Balance",      _fmt(cash_val),    "",    GREEN_RGB),
        ("Net Debt",          _fmt(debt_val),    "",    BLUE_RGB),
        ("FCF YTD",           _fmt(fcf_ytd),     "",    fcf_col),
        ("CAPEX Spent YTD",   _fmt(capex_act),   "",    YELLOW_RGB),
    ]
    _kpi_row_4(sl, kpis, top=CONTENT_TOP)

    fig_cash = sc.fig_cash_trend(data, height=380)
    if ar_data:
        gov     = ar_data.get("gov",     {})
        non_gov = ar_data.get("non_gov", {})
    else:
        gov = non_gov = {}
    fig_ar = sc.fig_ar_aging(gov, non_gov, height=340)

    _add_chart(sl, fig_cash,
               left=MARGIN_L, top=CHART_TOP_1,
               width=HALF_W, height=CHART_H_1,
               px_w=PX_HALF_W, px_h=PX_HALF_H)
    _add_chart(sl, fig_ar,
               left=MARGIN_L + HALF_W + GAP, top=CHART_TOP_1,
               width=HALF_W, height=CHART_H_1,
               px_w=PX_HALF_W, px_h=PX_HALF_H)


def _slide_consumer(prs: Presentation, data: dict, period: str) -> None:
    consumer = data.get("Consumer_Sales", pd.DataFrame())
    sl = _blank_slide(prs)
    _fill_bg(sl)
    _add_header(sl, "Consumer",
                "Prepaid · Postpaid · WTTx — Revenue & Subscribers", period)
    _add_footer(sl)

    rev_ytd  = _safe_sum(consumer, "Revenue")
    subs_val = float(consumer.groupby("Segment")["Subscribers"].last().sum()) \
               if not consumer.empty else 0.0
    arpu_val = float(consumer["ARPU"].mean()) if "ARPU" in consumer.columns else 0.0

    kpis = [
        ("Revenue YTD",    _fmt(rev_ytd),    "",    BLUE_RGB),
        ("Total Subs",     f"{subs_val:,.0f}", "",  GREEN_RGB),
        ("Blended ARPU",   f"{arpu_val:,.0f}", "Average", PURPLE_RGB),
    ]
    _kpi_row_3(sl, kpis, top=CONTENT_TOP)

    fig_rev  = sc.fig_consumer_revenue(data, height=380)
    fig_subs = sc.fig_consumer_subscribers(data, height=380)

    _add_chart(sl, fig_rev,
               left=MARGIN_L, top=CHART_TOP_1,
               width=HALF_W, height=CHART_H_1,
               px_w=PX_HALF_W, px_h=PX_HALF_H)
    _add_chart(sl, fig_subs,
               left=MARGIN_L + HALF_W + GAP, top=CHART_TOP_1,
               width=HALF_W, height=CHART_H_1,
               px_w=PX_HALF_W, px_h=PX_HALF_H)


def _slide_business_sales(prs: Presentation, data: dict, period: str) -> None:
    bs = data.get("Business_Sales", pd.DataFrame())
    sl = _blank_slide(prs)
    _fill_bg(sl)
    _add_header(sl, "Business Sales",
                "Revenue · Gross Profit · MRR by Segment", period)
    _add_footer(sl)

    rev_ytd = _safe_sum(bs, "Revenue")
    gp_ytd  = _safe_sum(bs, "Gross_Profit")
    mrr_ytd = _safe_sum(bs, "MRR")

    kpis = [
        ("Revenue YTD",     _fmt(rev_ytd), "",   BLUE_RGB),
        ("Gross Profit YTD", _fmt(gp_ytd), "",   GREEN_RGB),
        ("MRR YTD",          _fmt(mrr_ytd), "",  PURPLE_RGB),
    ]
    _kpi_row_3(sl, kpis, top=CONTENT_TOP)

    fig_rev = sc.fig_business_revenue(data, height=380)
    fig_gp  = sc.fig_business_gp(data, height=380)

    _add_chart(sl, fig_rev,
               left=MARGIN_L, top=CHART_TOP_1,
               width=HALF_W, height=CHART_H_1,
               px_w=PX_HALF_W, px_h=PX_HALF_H)
    _add_chart(sl, fig_gp,
               left=MARGIN_L + HALF_W + GAP, top=CHART_TOP_1,
               width=HALF_W, height=CHART_H_1,
               px_w=PX_HALF_W, px_h=PX_HALF_H)


def _slide_dpdi(prs: Presentation, data: dict, period: str) -> None:
    dpdi = data.get("DPDI", pd.DataFrame())
    sl = _blank_slide(prs)
    _fill_bg(sl)
    _add_header(sl, "DPDI — Digital Products",
                "e-Tender · e-GOVTT · e-Cashbook · e-Health · PAYPR", period)
    _add_footer(sl)

    rev_ytd = _safe_sum(dpdi, "Revenue")
    gp_ytd  = _safe_sum(dpdi, "Gross_Profit")
    ebitda_ytd = _safe_sum(dpdi, "EBITDA")

    kpis = [
        ("Revenue YTD",    _fmt(rev_ytd),    "",   BLUE_RGB),
        ("Gross Profit YTD", _fmt(gp_ytd),   "",   GREEN_RGB),
        ("EBITDA YTD",     _fmt(ebitda_ytd), "",   PURPLE_RGB),
    ]
    _kpi_row_3(sl, kpis, top=CONTENT_TOP)

    fig_prod  = sc.fig_dpdi_by_product(data, height=380)
    fig_trend = sc.fig_dpdi_monthly_trend(data, height=380)

    _add_chart(sl, fig_prod,
               left=MARGIN_L, top=CHART_TOP_1,
               width=HALF_W, height=CHART_H_1,
               px_w=PX_HALF_W, px_h=PX_HALF_H)
    _add_chart(sl, fig_trend,
               left=MARGIN_L + HALF_W + GAP, top=CHART_TOP_1,
               width=HALF_W, height=CHART_H_1,
               px_w=PX_HALF_W, px_h=PX_HALF_H)


def _slide_amplia(prs: Presentation, data: dict, period: str) -> None:
    amp = data.get("AMPLIA_Financial", pd.DataFrame())
    sl = _blank_slide(prs)
    _fill_bg(sl)
    _add_header(sl, "Amplia — Fibre Broadband Subsidiary",
                "Revenue · EBITDA · PAT", period)
    _add_footer(sl)

    rev_ytd    = _safe_sum(amp, "Revenue")
    ebitda_ytd = _safe_sum(amp, "EBITDA")
    pat_ytd    = _safe_sum(amp, "PAT")
    margin_v   = _safe_last(amp, "EBITDA_Margin") if "EBITDA_Margin" in amp.columns else None
    # compute margin from EBITDA/Revenue if not a direct column
    if margin_v is None and rev_ytd:
        margin_v = ebitda_ytd / rev_ytd * 100

    kpis = [
        ("Revenue YTD",    _fmt(rev_ytd),    "",   BLUE_RGB),
        ("EBITDA YTD",     _fmt(ebitda_ytd), "",   GREEN_RGB),
        ("PAT YTD",        _fmt(pat_ytd),    "",   PURPLE_RGB),
        ("EBITDA Margin",  f"{margin_v:.1f}%" if margin_v else "—", "YTD avg", YELLOW_RGB),
    ]
    _kpi_row_4(sl, kpis, top=CONTENT_TOP)

    fig_rev   = sc.fig_amplia_revenue(data, height=380)
    fig_ebitda = sc.fig_amplia_ebitda(data, height=380)
    fig_pat   = sc.fig_amplia_pat(data, height=380)

    chart_h = CHART_H_1

    _add_chart(sl, fig_rev,
               left=MARGIN_L, top=CHART_TOP_1,
               width=THIRD_W, height=chart_h,
               px_w=PX_THIRD_W, px_h=PX_THIRD_H)
    _add_chart(sl, fig_ebitda,
               left=MARGIN_L + THIRD_W + GAP, top=CHART_TOP_1,
               width=THIRD_W, height=chart_h,
               px_w=PX_THIRD_W, px_h=PX_THIRD_H)
    _add_chart(sl, fig_pat,
               left=MARGIN_L + 2 * (THIRD_W + GAP), top=CHART_TOP_1,
               width=THIRD_W, height=chart_h,
               px_w=PX_THIRD_W, px_h=PX_THIRD_H)


# ══════════════════════════════════════════════════════════════════════════════
# PUBLIC ENTRY POINT
# ══════════════════════════════════════════════════════════════════════════════

def build_pptx(data: dict, ar_data: Optional[dict] = None) -> io.BytesIO:
    """Build the full TSTT board pack.

    Args:
        data:    dict from load_all_data()
        ar_data: dict from load_ar_aging(), or None for fallback values

    Returns:
        BytesIO containing the .pptx file, seeked to position 0.
    """
    period = _period(data)
    prs = _new_prs()

    _slide_title(prs, period)
    _slide_scorecard(prs, data, period)
    _slide_financial(prs, data, period)
    _slide_ebitda_bridge(prs, data, period)
    _slide_opex(prs, data, period)
    _slide_cash_capex(prs, data, ar_data, period)
    _slide_consumer(prs, data, period)
    _slide_business_sales(prs, data, period)
    _slide_dpdi(prs, data, period)
    _slide_amplia(prs, data, period)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
