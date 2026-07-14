"""
TSTT Board Report — high-fidelity screenshot PowerPoint exporter.

Drives headless Chromium over the live Streamlit dashboard, screenshots every
page (Home + each sidebar page), and assembles the captures into a 16:9 .pptx.
Pages taller than a slide are auto-split into multiple legible slides.

Usage (from the project root):
    python tools/pptx_screenshot_export.py
    python tools/pptx_screenshot_export.py --url http://localhost:8501
    python tools/pptx_screenshot_export.py --page Consumer_Sales

One-time setup:
    pip install playwright Pillow
    playwright install chromium
"""

from __future__ import annotations

import argparse
import io
import re
import sys
import time
import urllib.parse
import urllib.request
from pathlib import Path

# ── Geometry ─────────────────────────────────────────────────────────────────
SLIDE_W_IN   = 13.333
SLIDE_H_IN   = 7.5
SLIDE_ASPECT = SLIDE_W_IN / SLIDE_H_IN          # 1.7778
BG_RGB       = (0, 0, 0)                        # app background (#000000)
OVERLAP_FRAC = 0.06                             # slice overlap so cards aren't cut

PROJECT_ROOT = Path(__file__).resolve().parent.parent

# CSS injected to strip Streamlit chrome from the screenshot.
HIDE_CSS = """
section[data-testid="stSidebar"]            { display: none !important; }
div[data-testid="stSidebarCollapsedControl"]{ display: none !important; }
button[data-testid="stSidebarCollapseButton"]{ display: none !important; }
[data-testid="stToolbar"]                   { display: none !important; }
[data-testid="stDecoration"]                { display: none !important; }
[data-testid="stHeader"]                    { display: none !important; }
header[data-testid="stHeader"]              { display: none !important; }
.stAppDeployButton                          { display: none !important; }
[data-testid="stStatusWidget"]              { display: none !important; }
"""


# ── Page discovery ───────────────────────────────────────────────────────────
def discover_pages() -> list[tuple[str, str]]:
    """Return [(label, slug), ...] for the Home page plus every pages/*.py file,
    in sidebar (numeric-prefix) order. Slug "" is the Home page; the Export page
    is skipped. The slug is the filename with its `<n>_` prefix and `.py` removed,
    matching Streamlit's multipage URL routing."""
    items = []
    for f in (PROJECT_ROOT / "pages").glob("*.py"):
        m = re.match(r"(\d+)_(.+)\.py$", f.name)
        if not m:
            continue
        num, slug = int(m.group(1)), m.group(2)
        if slug == "Export":
            continue
        items.append((num, slug))
    items.sort(key=lambda t: t[0])
    return [("Home", "")] + [(slug.replace("_", " "), slug) for _, slug in items]


# ── Streamlit server ─────────────────────────────────────────────────────────
def _wait_for_server(url: str, timeout: float = 120.0) -> None:
    deadline = time.time() + timeout
    last_err: Exception | None = None
    while time.time() < deadline:
        try:
            with urllib.request.urlopen(url, timeout=3) as resp:
                if resp.status == 200:
                    return
        except Exception as exc:               # noqa: BLE001 — any failure = not ready
            last_err = exc
        time.sleep(1.0)
    raise RuntimeError(
        f"Streamlit server never became ready at {url} within {timeout}s "
        f"(last error: {last_err})"
    )


def _launch_streamlit(port: int):
    import subprocess

    cmd = [
        sys.executable, "-m", "streamlit", "run", str(PROJECT_ROOT / "app.py"),
        "--server.port", str(port),
        "--server.headless", "true",
        "--server.runOnSave", "false",
        "--browser.gatherUsageStats", "false",
    ]
    print(f"  launching Streamlit on port {port} ...")
    return subprocess.Popen(
        cmd, cwd=str(PROJECT_ROOT),
        stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
    )


# ── Capture ──────────────────────────────────────────────────────────────────
def _shoot(page) -> bytes:
    """Settle the page (charts + below-the-fold) and return a full-page PNG."""
    page.evaluate("window.scrollTo(0, document.body.scrollHeight)")
    page.wait_for_timeout(900)
    page.evaluate("window.scrollTo(0, 0)")
    page.wait_for_timeout(500)
    return page.screenshot(full_page=True)


def _capture_page(page, base_url: str, label: str, slug: str,
                  focus_month: str | None = None):
    """Capture one dashboard page; returns [(caption, png_bytes), ...].

    A page with st.tabs yields one capture per tab; a tab-less page yields one."""
    root = base_url.rstrip("/")
    url  = f"{root}/{slug}?embed=true" if slug else f"{root}/?embed=true"
    if focus_month:
        url += f"&focus_month={urllib.parse.quote(focus_month)}"
    print(f"  opening {label}  ->  {url}")
    page.goto(url, wait_until="domcontentloaded", timeout=60_000)
    page.add_style_tag(content=HIDE_CSS)            # fresh document each goto
    try:
        page.wait_for_selector('div[data-testid="stAppViewContainer"]', timeout=45_000)
    except Exception:                               # noqa: BLE001
        pass
    page.wait_for_timeout(3500)                     # Streamlit websocket render + Plotly

    tabs   = page.locator('button[data-baseweb="tab"]')
    n_tabs = tabs.count()
    out: list[tuple[str, bytes]] = []

    if n_tabs == 0:
        out.append((label, _shoot(page)))
        print(f"  captured {label}")
    else:
        for i in range(n_tabs):
            tab = tabs.nth(i)
            tab_label = (tab.inner_text() or f"Tab {i + 1}").strip()
            tab.click()
            page.evaluate("window.dispatchEvent(new Event('resize'))")
            page.wait_for_timeout(2500)             # settle: Plotly animation + layout
            out.append((f"{label} — {tab_label}", _shoot(page)))
            print(f"  captured {label} — {tab_label}")
    return out


def capture_all(base_url: str, pages: list[tuple[str, str]], viewport_w: int,
                headed: bool, focus_month: str | None = None):
    """Capture every page in `pages` with a single shared browser session."""
    from playwright.sync_api import sync_playwright

    captures: list[tuple[str, bytes]] = []
    with sync_playwright() as pw:
        try:
            browser = pw.chromium.launch(headless=not headed)
        except Exception as exc:                   # noqa: BLE001
            raise RuntimeError(
                "Could not launch Chromium. Run `playwright install chromium` once. "
                f"Original error: {exc}"
            ) from exc

        context = browser.new_context(
            viewport={"width": viewport_w, "height": round(viewport_w / SLIDE_ASPECT)},
            device_scale_factor=2,                  # 2x density => crisp text
        )
        page = context.new_page()

        for label, slug in pages:
            try:
                captures += _capture_page(page, base_url, label, slug, focus_month)
            except Exception as exc:               # noqa: BLE001 — skip, keep going
                print(f"  WARNING: failed to capture {label}: {exc}", file=sys.stderr)

        context.close()
        browser.close()

    return captures


# ── Slicing ──────────────────────────────────────────────────────────────────
def slice_capture(png_bytes: bytes, split: bool):
    """Slice one full-page screenshot into 16:9 segments (PIL Images)."""
    from PIL import Image

    img = Image.open(io.BytesIO(png_bytes)).convert("RGB")
    w, h = img.size
    seg_h = round(w / SLIDE_ASPECT)

    # Shorter than a slide, or splitting disabled => single frame.
    if not split or h <= seg_h:
        if h >= seg_h:
            return [img]                            # caller fits it onto the slide
        canvas = Image.new("RGB", (w, seg_h), BG_RGB)
        canvas.paste(img, (0, 0))
        return [canvas]

    overlap = round(seg_h * OVERLAP_FRAC)
    step    = seg_h - overlap
    segments, top = [], 0
    while True:
        if top + seg_h >= h:                        # last slice anchored to the bottom
            segments.append(img.crop((0, h - seg_h, w, h)))
            break
        segments.append(img.crop((0, top, w, top + seg_h)))
        top += step
    return segments


# ── PPTX assembly ────────────────────────────────────────────────────────────
def build_pptx(captures, out_path: Path, split: bool) -> int:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    slide_w, slide_h = prs.slide_width, prs.slide_height
    n_slides = 0

    for label, png in captures:
        for seg in slice_capture(png, split):
            slide = prs.slides.add_slide(blank)
            bg = slide.background
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(*BG_RGB)

            bio = io.BytesIO()
            seg.save(bio, format="PNG")
            bio.seek(0)

            sw, sh = seg.size
            if split:
                # Segments are already exact 16:9 => place full-bleed.
                slide.shapes.add_picture(bio, 0, 0, width=slide_w, height=slide_h)
            else:
                # Fit whole capture onto one slide, preserve aspect, centre it.
                if sw / sh >= SLIDE_ASPECT:
                    pic_w, pic_h = slide_w, int(slide_w * sh / sw)
                else:
                    pic_h, pic_w = slide_h, int(slide_h * sw / sh)
                slide.shapes.add_picture(
                    bio, int((slide_w - pic_w) / 2), int((slide_h - pic_h) / 2),
                    width=pic_w, height=pic_h,
                )
            n_slides += 1

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"\n  wrote {n_slides} slide(s) -> {out_path}")
    return n_slides


# ── Main ─────────────────────────────────────────────────────────────────────
def main() -> int:
    ap = argparse.ArgumentParser(description="Screenshot-based all-pages PPTX exporter.")
    ap.add_argument("--url", default=None,
                    help="Reuse a running Streamlit server (e.g. http://localhost:8501). "
                         "If omitted, a dedicated instance is launched.")
    ap.add_argument("--port", type=int, default=8599,
                    help="Port for the auto-launched Streamlit instance (default 8599).")
    ap.add_argument("--page", default=None,
                    help="Optional: capture only this one page slug (default: all pages).")
    ap.add_argument("--pages", default=None,
                    help="Optional: comma-separated page labels/slugs to capture "
                         "(default: all). Takes precedence over --page.")
    ap.add_argument("--focus-month", default=None,
                    help="Focus month to render every page at (e.g. Jun-26). "
                         "Passed to the app as a query param; default: latest month.")
    ap.add_argument("--out", default="exports/TSTT_Board_Report.pptx",
                    help="Output .pptx path (default: exports/TSTT_Board_Report.pptx).")
    ap.add_argument("--viewport-width", type=int, default=1920,
                    help="Browser viewport width in CSS px (default 1920).")
    ap.add_argument("--no-split", action="store_true",
                    help="One slide per page/tab, scaled to fit (no auto-split).")
    ap.add_argument("--headed", action="store_true",
                    help="Show the browser window (debugging).")
    args = ap.parse_args()

    out_path = Path(args.out)
    if not out_path.is_absolute():
        out_path = PROJECT_ROOT / out_path

    pages = discover_pages()
    if args.pages:
        wanted = {p.strip() for p in args.pages.split(",") if p.strip()}
        pages = [(l, s) for l, s in pages if l in wanted or s in wanted]
        if not pages:
            print(f"ERROR: no pages match --pages {args.pages!r}.", file=sys.stderr)
            return 1
    elif args.page:
        pages = [(l, s) for l, s in pages if args.page in (s, l)]
        if not pages:
            print(f"ERROR: no page matches --page {args.page!r}.", file=sys.stderr)
            return 1
    print(f"Pages to capture: {', '.join(l for l, _ in pages)}")

    proc = None
    try:
        if args.url:
            base_url = args.url
            print(f"Using existing server: {base_url}")
            _wait_for_server(base_url)
        else:
            proc = _launch_streamlit(args.port)
            base_url = f"http://localhost:{args.port}"
            _wait_for_server(base_url)
            print("  server ready")

        captures = capture_all(base_url, pages, args.viewport_width, args.headed,
                               args.focus_month)
        if not captures:
            print("ERROR: no pages captured.", file=sys.stderr)
            return 1

        build_pptx(captures, out_path, split=not args.no_split)
        print("Done.")
        return 0

    finally:
        if proc is not None:
            print("  shutting down Streamlit ...")
            proc.terminate()
            try:
                proc.wait(timeout=10)
            except Exception:                       # noqa: BLE001
                proc.kill()


if __name__ == "__main__":
    raise SystemExit(main())
